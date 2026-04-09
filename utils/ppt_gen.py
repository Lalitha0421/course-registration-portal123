from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from io import BytesIO
from datetime import datetime
from config import get_connection

def generate_registration_ppt(student_id, current_session):
    conn = get_connection()
    cursor = conn.cursor()
    
    try:
        # Get student info
        cursor.execute("""
            SELECT name, enrollment_no, program, email, mobile, semester
            FROM students WHERE student_id = :sid
        """, sid=student_id)
        student = cursor.fetchone()
        if not student:
            return None
        
        name, enrl, prog, email, mob, sem = student
        
        # Get registered courses
        cursor.execute("""
            SELECT cm.course_code, cm.course_title, cm.course_type, cm.credits,
                   f.faculty_name
            FROM student_registration_courses src
            JOIN course_instance ci ON src.course_instance_id = ci.instance_id
            JOIN course_master cm ON ci.course_id = cm.course_id
            LEFT JOIN faculty f ON ci.faculty_id = f.faculty_id
            JOIN student_registration sr ON src.reg_id = sr.reg_id
            WHERE sr.student_id = :sid AND sr.academic_session = :sess
        """, sid=student_id, sess=current_session)
        courses = cursor.fetchall()

        # Create PPT
        prs = Presentation()
        
        # Slide 1: Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Course Registration Report"
        subtitle.text = f"Student: {name}\nRegistration Session: {current_session}\nVNIT Nagpur"
        
        # Slide 2: Student Information
        info_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(info_slide_layout)
        title = slide.shapes.title
        title.text = "Student Information"
        
        tf = slide.placeholders[1].text_frame
        tf.text = f"Name: {name}"
        p = tf.add_paragraph()
        p.text = f"Enrollment No: {enrl}"
        p = tf.add_paragraph()
        p.text = f"Program: {prog}"
        p = tf.add_paragraph()
        p.text = f"Semester: {sem}"
        p = tf.add_paragraph()
        p.text = f"Contact: {mob} | {email}"
        
        # Slide 3: Registered Courses
        course_slide_layout = prs.slide_layouts[5] # Blank with Title
        slide = prs.slides.add_slide(course_slide_layout)
        title = slide.shapes.title
        title.text = "Registered Courses Summary"
        
        # Add Table
        rows = len(courses) + 1
        cols = 4
        left = Inches(1.0)
        top = Inches(1.5)
        width = Inches(8.0)
        height = Inches(0.8)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Header
        table.cell(0, 0).text = "Code"
        table.cell(0, 1).text = "Title"
        table.cell(0, 2).text = "Type"
        table.cell(0, 3).text = "Credits"
        
        for i, row in enumerate(courses, 1):
            table.cell(i, 0).text = str(row[0])
            table.cell(i, 1).text = str(row[1])
            table.cell(i, 2).text = str(row[2])
            table.cell(i, 3).text = str(row[3])
            
        # Slide 4: Registration Summary
        summary_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(summary_slide_layout)
        title = slide.shapes.title
        title.text = "Registration Summary"
        
        total_credits = sum(r[3] for r in courses)
        dc_count = sum(1 for r in courses if r[2] == 'DC')
        de_count = sum(1 for r in courses if r[2] == 'DE')
        
        tf = slide.placeholders[1].text_frame
        tf.text = f"Total Courses: {len(courses)}"
        p = tf.add_paragraph()
        p.text = f"Total Credits: {total_credits}"
        p = tf.add_paragraph()
        p.text = f"Departmental Core (DC): {dc_count}"
        p = tf.add_paragraph()
        p.text = f"Departmental Electives (DE): {de_count}"
        p = tf.add_paragraph()
        p.text = f"Registration Status: Submitted / Awaiting Approval"
        
        # Footer Slide
        footer_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(footer_slide_layout)
        title = slide.shapes.title
        title.text = "End of Report"
        slide.placeholders[1].text = f"Generated Date: {datetime.now().strftime('%Y-%m-%d %H:%M')}\nVNIT Academic Portal"

        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output
        
    except Exception as e:
        print(f"PPT generation failed: {str(e)}")
        return None
    finally:
        cursor.close()
        conn.close()
